"""
CVMatchMaker Presentation Generator — Updated Version
Generates a clean 10-slide PowerPoint with new screenshots
and the 3 key technical points (Content Extraction, NLP, Recommendations).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SHOTS = "/Users/saifkhan/Desktop/Final year/app/scripts/shots"
OUTPUT = "/Users/saifkhan/Desktop/Final year/CVMatchMaker_Presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colours
BG    = RGBColor(17, 17, 17)
GREEN = RGBColor(34, 197, 94)
WHITE = RGBColor(240, 240, 240)
GREY  = RGBColor(150, 150, 150)
FONT  = "Avenir Next"


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def bar(slide):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()


def txt(slide, l, t, w, h, text, sz=18, bold=False, col=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text; p.alignment = align
    for r in p.runs:
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT


def body(slide, l, t, w, h, lines, sz=16):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        p = box.text_frame.paragraphs[0] if i == 0 else box.text_frame.add_paragraph()
        p.text = line; p.space_after = Pt(8)
        p.font.size = Pt(sz); p.font.color.rgb = WHITE; p.font.name = FONT


def img(slide, name, l, t, w=None):
    path = f"{SHOTS}/{name}"
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, width=w)


# ── SLIDE 1: Title ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
s.shapes.add_shape(1, Inches(0), Inches(3.2), prs.slide_width, Inches(0.06)).fill.solid()
s.shapes[-1].fill.fore_color.rgb = GREEN; s.shapes[-1].line.fill.background()
txt(s, Inches(1), Inches(1.8), Inches(11), Inches(1), "CVMatchMaker", 52, True, GREEN, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.6), Inches(11), Inches(0.5), "AI-Powered CV Analysis Tool", 22, False, GREY, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.8), Inches(11), Inches(1),
    "Saif Ali Khan  ·  BEng Computer Science  ·  University of Roehampton  ·  2025", 16, False, GREY, PP_ALIGN.CENTER)

# ── SLIDE 2: Problem & Solution ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "Problem & Solution", 34, True, GREEN)
body(s, Inches(0.8), Inches(1.3), Inches(11), Inches(5.5), [
    "The Problem:",
    "When you apply for a job, your CV goes through an ATS (Applicant Tracking System) — a computer that scans for keywords. If the right skills are missing, your CV gets rejected before any human sees it. Most CV tools only check spelling and grammar — they don't compare your CV to real jobs.",
    "",
    "The Solution:",
    "CVMatchMaker uses NLP techniques to read your CV and compare it against real job listings from Reed.co.uk. It tells you exactly which skills are missing, gives you a score out of 100, and links you to free courses so you can improve.",
    "",
    "How it works:  Upload CV → Extract content → NLP keyword grabbing → Compare with live jobs → Show results & recommendations"
], 18)

# ── SLIDE 3: Technical Pipeline (NEW — the 3 key points) ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "How It Works — 3 Step Pipeline", 34, True, GREEN)
body(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), [
    "Step 1 — Content Extraction",
    "Using content extractor to extract CV content from PDF and DOCX.",
    "• cv_parser.py uses PyPDF2 to read PDFs and python-docx to read Word files",
    "• Extracts raw text and splits it into sections (Skills, Experience, Education)",
    "",
    "Step 2 — NLP Keyword Analysis",
    "Apply natural language processing techniques to analyse CV content, like grab keywords after processing.",
    "• skill_extractor.py uses spaCy (NLP library) with 3 layers:",
    "  Layer 1: Regex phrase matching (finds exact skill names)",
    "  Layer 2: Lemmatisation (\"developing\" → \"develop\")",
    "  Layer 3: Synonym normalisation (\"reactjs\" → \"react\")",
], 16)
body(s, Inches(6.8), Inches(1.2), Inches(5.5), Inches(5.5), [
    "Step 3 — Recommendations",
    "Using recommendation/search techniques to provide proper recommendations based on the CV.",
    "• job_fetcher.py connects to Reed.co.uk API for real job listings",
    "• analyser.py compares CV skills vs job skills using set operations",
    "• Missing skills are ranked by frequency (how often employers ask for them)",
    "• Each missing skill is mapped to free learning resources",
    "",
    "The Pipeline Flow:",
    "Upload CV → cv_parser.py → skill_extractor.py → job_fetcher.py → analyser.py → Results Page",
], 16)

# ── SLIDE 4: Homepage ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.6), "Homepage", 30, True, GREEN)
body(s, Inches(0.5), Inches(1.0), Inches(4.3), Inches(5.5), [
    "This is the first page users see.",
    "",
    "It shows a clear message: 'Upload Your CV. See What's Missing.'",
    "",
    "Below it are Login and Sign Up forms.",
    "",
    "Passwords are stored securely using hashing (PBKDF2-SHA256) — they can never be read by anyone.",
], 16)
img(s, "01_homepage.png", Inches(5.3), Inches(0.4), w=Inches(7.7))

# ── SLIDE 5: Dashboard ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.6), "Dashboard", 30, True, GREEN)
body(s, Inches(0.5), Inches(1.0), Inches(4.3), Inches(5.5), [
    "After logging in, the user sees their Dashboard.",
    "",
    "They upload a CV file (PDF or Word), type the job title they want, and click Analyse.",
    "",
    "The system reads the CV using content extraction, calls the Reed API for live jobs, runs NLP skill extraction, and generates results — all in a few seconds.",
], 16)
img(s, "02_dashboard.png", Inches(5.3), Inches(0.4), w=Inches(7.7))

# ── SLIDE 6: Results — Score & Skills ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.6), "Results — Score & Skills", 30, True, GREEN)
body(s, Inches(0.5), Inches(1.0), Inches(4.3), Inches(5.5), [
    "The Career Score (out of 100) shows how well the CV matches real job listings.",
    "",
    "Green chips = skills found on the CV (matched using NLP).",
    "Red chips = skills that are missing.",
    "",
    "Each missing skill has a priority tag (High/Medium/Low) based on how many employers ask for it.",
    "",
    "Links to free courses on YouTube, Coursera, and W3Schools are provided for each missing skill.",
], 16)
img(s, "04_results_top.png", Inches(5.3), Inches(0.4), w=Inches(7.7))

# ── SLIDE 7: History & Compare ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.6), "History & Compare", 30, True, GREEN)
body(s, Inches(0.5), Inches(1.0), Inches(4.3), Inches(5.5), [
    "History: Every analysis is saved. A Score Trend chart shows how the user is improving over time.",
    "",
    "Compare: Users can pick two past analyses and see them side-by-side — scores, matched skills, and new skills gained.",
    "",
    "Users can also export reports and delete old analyses.",
], 16)
img(s, "03_history.png", Inches(5.3), Inches(0.4), w=Inches(7.7))

# ── SLIDE 8: Learning Path ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.6), "Learning Path", 30, True, GREEN)
body(s, Inches(0.5), Inches(1.0), Inches(4.3), Inches(5.5), [
    "Users can add missing skills to their personal Learning Path — like a to-do list for skills.",
    "",
    "Each skill has links to free courses. When they finish learning, they mark it as done.",
    "",
    "A progress bar shows how many skills they've completed. All data is saved in the database.",
], 16)
img(s, "08_learning.png", Inches(5.3), Inches(0.4), w=Inches(7.7))

# ── SLIDE 9: Tech Stack ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "How It Was Built", 34, True, GREEN)
body(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Backend: Python + Flask",
    "Handles all the logic — reading CVs, calling the API, running the NLP analysis, and serving web pages.",
    "",
    "NLP Engine: spaCy",
    "A Natural Language Processing library that reads text, tokenises words, and extracts skill names using lemmatisation and pattern matching.",
    "",
    "Job Data: Reed.co.uk API",
    "Downloads real UK job listings so the analysis is always based on current market data.",
], 17)
body(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Database: SQLite (3 tables)",
    "Users, AnalysisSessions, LearningGoals. CV files are NEVER saved — only scores and skills.",
    "",
    "Frontend: HTML, CSS, JavaScript",
    "Clean responsive pages. No complex framework — just standard web technologies (W3Schools style).",
    "",
    "Security: PBKDF2 password hashing, Flask-Login sessions, GDPR compliant (CVs processed in memory only and deleted immediately).",
], 17)

# ── SLIDE 10: Thank You ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
s.shapes.add_shape(1, Inches(0), Inches(3.2), prs.slide_width, Inches(0.06)).fill.solid()
s.shapes[-1].fill.fore_color.rgb = GREEN; s.shapes[-1].line.fill.background()
txt(s, Inches(1), Inches(2.0), Inches(11), Inches(1), "Thank You", 52, True, GREEN, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.0), Inches(11), Inches(0.5), "Any Questions?", 22, False, GREY, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.2), Inches(11), Inches(1),
    "Saif Ali Khan  ·  University of Roehampton  ·  2025\nPython Flask  ·  spaCy NLP  ·  Reed API  ·  SQLite",
    16, False, GREY, PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved {len(prs.slides)} slides to {OUTPUT}")
